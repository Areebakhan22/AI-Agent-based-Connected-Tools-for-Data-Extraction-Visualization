"""
PowerPoint Generator Module using python-pptx

This module creates PowerPoint (.pptx) presentations from structured SysML data.
Uses the same layout calculation logic as Google Slides generator for consistency.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import Dict, List
import os
from slides_generator import (
    calculate_professional_layout,
    SLIDE_WIDTH,
    SLIDE_HEIGHT
)

# PowerPoint uses different units - convert from points to inches
# 1 inch = 72 points
PT_TO_INCH = 1.0 / 72.0


def pt_to_inches(points: float) -> float:
    """Convert points to inches for PowerPoint."""
    return points * PT_TO_INCH


def generate_pptx(data: Dict, output_filename: str = None, title: str = "SysML Visualization") -> str:
    """
    Generate a PowerPoint presentation from SysML data.
    
    Args:
        data: Dictionary with parts, actors, use_cases, connections, hierarchy
        output_filename: Output filename (default: based on system name)
        title: Presentation title
        
    Returns:
        Path to generated .pptx file
    """
    # Calculate layout (reuse same logic as Google Slides)
    layout = calculate_professional_layout(data)
    
    # Create presentation
    prs = Presentation()
    # Set slide size to match Google Slides (16:9 aspect ratio)
    prs.slide_width = Inches(10)  # 10 inches = 720 points
    prs.slide_height = Inches(7.5)  # 7.5 inches = 540 points
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Draw system boundary
    boundary = layout['system_boundary']
    boundary_shape = draw_system_boundary_pptx(
        slide,
        boundary['name'],
        pt_to_inches(boundary['x']),
        pt_to_inches(boundary['y']),
        pt_to_inches(boundary['width']),
        pt_to_inches(boundary['height'])
    )
    
    # Store shape IDs for connections
    shape_map = {}
    
    # Draw all elements
    for elem_name, elem_info in layout['elements'].items():
        elem_type = elem_info['type']
        
        if elem_type == 'part':
            # Parts use top-left coordinates
            x = pt_to_inches(elem_info['x'])
            y = pt_to_inches(elem_info['y'])
            shape = draw_part_pptx(
                slide,
                elem_name,
                x,
                y,
                pt_to_inches(elem_info['width']),
                pt_to_inches(elem_info['height'])
            )
            shape_map[elem_name] = shape
            
        elif elem_type == 'use_case':
            # Use cases use top-left coordinates
            x = pt_to_inches(elem_info['x'])
            y = pt_to_inches(elem_info['y'])
            shape = draw_use_case_pptx(
                slide,
                elem_name,
                x,
                y,
                pt_to_inches(elem_info['width']),
                pt_to_inches(elem_info['height'])
            )
            shape_map[elem_name] = shape
            
        elif elem_type == 'actor':
            # Actors use center coordinates, convert to top-left for drawing
            center_x = pt_to_inches(elem_info['x'])
            center_y = pt_to_inches(elem_info['y'])
            size = pt_to_inches(elem_info['size'])
            x = center_x - size / 2
            y = center_y - size / 2
            shape = draw_actor_pptx(
                slide,
                elem_name,
                center_x,  # Pass center for connection calculations
                center_y,
                size
            )
            shape_map[elem_name] = shape
    
    # Draw connections
    connections = data.get('connections', [])
    for conn in connections:
        from_name = conn['from']
        to_name = conn['to']
        
        if from_name in shape_map and to_name in shape_map:
            from_shape = shape_map[from_name]
            to_shape = shape_map[to_name]
            from_layout = layout['elements'][from_name]
            to_layout = layout['elements'][to_name]
            
            # Convert layout to proper format for connection drawing
            # For actors, we need to convert center+size to bounding box
            from_layout_for_conn = from_layout.copy()
            if from_layout['type'] == 'actor':
                from_layout_for_conn['x'] = from_layout['x'] - from_layout['size'] / 2
                from_layout_for_conn['y'] = from_layout['y'] - from_layout['size'] / 2
                from_layout_for_conn['width'] = from_layout['size']
                from_layout_for_conn['height'] = from_layout['size']
            
            to_layout_for_conn = to_layout.copy()
            if to_layout['type'] == 'actor':
                to_layout_for_conn['x'] = to_layout['x'] - to_layout['size'] / 2
                to_layout_for_conn['y'] = to_layout['y'] - to_layout['size'] / 2
                to_layout_for_conn['width'] = to_layout['size']
                to_layout_for_conn['height'] = to_layout['size']
            
            draw_connection_pptx(
                slide,
                from_shape,
                to_shape,
                from_layout_for_conn,
                to_layout_for_conn
            )
    
    # Determine output filename
    if not output_filename:
        system_name = boundary['name']
        safe_name = "".join(c for c in system_name if c.isalnum() or c in (' ', '-', '_')).strip()
        output_filename = f"{safe_name}_SysML.pptx"
    
    # Ensure .pptx extension
    if not output_filename.endswith('.pptx'):
        output_filename += '.pptx'
    
    # Save presentation
    prs.save(output_filename)
    return os.path.abspath(output_filename)


def draw_system_boundary_pptx(slide, system_name: str, x: float, y: float, 
                               width: float, height: float):
    """
    Draw system boundary as a large rectangle with title.
    """
    # Create rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height)
    )
    
    # Style: Light gray fill, dark border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
    
    line = shape.line
    line.color.rgb = RGBColor(100, 100, 100)  # Dark gray
    line.width = Pt(2)
    
    # Add title text in top-left
    text_frame = shape.text_frame
    text_frame.text = system_name
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    
    # Format title text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0, 0, 0)
    
    return shape


def draw_part_pptx(slide, part_name: str, x: float, y: float, 
                   width: float, height: float):
    """
    Draw a part as a rectangle.
    """
    # Adjust x, y to be top-left corner (layout provides center for actors)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height)
    )
    
    # Style: White fill, dark border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    line = shape.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Pt(1.5)
    
    # Add text
    text_frame = shape.text_frame
    text_frame.text = part_name
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Format text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0, 0, 0)
    
    return shape


def draw_use_case_pptx(slide, use_case_name: str, x: float, y: float,
                       width: float, height: float):
    """
    Draw a use case as a rounded rectangle.
    """
    # Adjust x, y to be top-left corner
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height)
    )
    
    # Adjust corner radius (percentage of width)
    shape.adjustments[0] = 0.1  # 10% corner radius
    
    # Style: Light blue fill, dark border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 240, 255)  # Light blue
    
    line = shape.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Pt(1.5)
    
    # Add text
    text_frame = shape.text_frame
    text_frame.text = use_case_name
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Format text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0, 0, 0)
    
    return shape


def draw_actor_pptx(slide, actor_name: str, center_x: float, center_y: float, size: float):
    """
    Draw an actor as a circle.
    
    Args:
        center_x, center_y: Center coordinates (in inches)
        size: Diameter (in inches)
    """
    # Convert center to top-left
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(center_x - size/2),
        Inches(center_y - size/2),
        Inches(size),
        Inches(size)
    )
    
    # Style: Light yellow fill, dark border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 200)  # Light yellow
    
    line = shape.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Pt(1.5)
    
    # Add text
    text_frame = shape.text_frame
    text_frame.text = actor_name
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Format text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0, 0, 0)
    
    return shape


def draw_connection_pptx(slide, from_shape, to_shape, from_layout: Dict, to_layout: Dict):
    """
    Draw a connection line between two shapes with arrowhead.
    
    Args:
        from_layout, to_layout: Layout dictionaries with x, y, width, height (top-left coordinates)
    """
    # Both layouts now use top-left coordinates with width/height
    from_left = from_layout['x']
    from_top = from_layout['y']
    from_width = from_layout['width']
    from_height = from_layout['height']
    from_right = from_left + from_width
    from_bottom = from_top + from_height
    from_center_x = from_left + from_width / 2
    from_center_y = from_top + from_height / 2
    
    to_left = to_layout['x']
    to_top = to_layout['y']
    to_width = to_layout['width']
    to_height = to_layout['height']
    to_right = to_left + to_width
    to_bottom = to_top + to_height
    to_center_x = to_left + to_width / 2
    to_center_y = to_top + to_height / 2
    
    # Calculate direction
    dx = to_center_x - from_center_x
    dy = to_center_y - from_center_y
    
    # Find intersection points with shape boundaries
    if abs(dx) > abs(dy):
        # Horizontal connection
        if dx > 0:
            start_x = from_right
            start_y = from_center_y
            end_x = to_left
            end_y = to_center_y
        else:
            start_x = from_left
            start_y = from_center_y
            end_x = to_right
            end_y = to_center_y
    else:
        # Vertical connection
        if dy > 0:
            start_x = from_center_x
            start_y = from_bottom
            end_x = to_center_x
            end_y = to_top
        else:
            start_x = from_center_x
            start_y = from_top
            end_x = to_center_x
            end_y = to_bottom
    
    # Create connector line (convert points to inches)
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(pt_to_inches(start_x)),
        Inches(pt_to_inches(start_y)),
        Inches(pt_to_inches(end_x)),
        Inches(pt_to_inches(end_y))
    )
    
    # Style the connector
    line = connector.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Pt(1.5)
    
    # Add arrowhead at end (if available in this version)
    try:
        # Try to set arrowhead properties (may not be available in all versions)
        if hasattr(line, 'end_arrowhead_length'):
            line.end_arrowhead_length = 1  # Short
        if hasattr(line, 'end_arrowhead_style'):
            line.end_arrowhead_style = 2  # Triangle
    except:
        pass  # Arrowheads not supported in this version
    
    return connector

